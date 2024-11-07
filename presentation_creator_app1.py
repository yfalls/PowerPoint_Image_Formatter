
# to run locally on terminal use: python -m streamlit run presentation_creator_app1.py

'''
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image
from zipfile import ZipFile
import os
import tempfile

def extract_images_from_zip(zip_file, extract_to):
    image_extensions = {".jpg", ".jpeg", ".png", ".bmp", ".gif"}
    image_paths = []
    with ZipFile(zip_file) as zip_ref:
        # Extract all files and gather only valid image files, excluding macOS hidden files
        for file_info in zip_ref.infolist():
            if not file_info.is_dir():
                # Check for valid image extensions and exclude macOS hidden files
                if (os.path.splitext(file_info.filename)[1].lower() in image_extensions and 
                    not os.path.basename(file_info.filename).startswith("._")):
                    extracted_path = zip_ref.extract(file_info, path=extract_to)
                    image_paths.append(extracted_path)
    return image_paths

def create_widescreen_presentation_with_images(output_path, images, crop_left=0, crop_bottom=0):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    if not images:
        st.warning("No images found to add to the presentation.")
        return

    for image in images:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        try:
            with Image.open(image) as img:
                original_width, original_height = img.size
                crop_box = (crop_left, 0, original_width, original_height - crop_bottom)
                cropped_img = img.crop(crop_box)
                
                temp_image_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                cropped_img.save(temp_image_path.name)
                
                img_shape = slide.shapes.add_picture(temp_image_path.name, 0, 0)
                aspect_ratio = img_shape.width / img_shape.height
                new_height = Inches(6.00)
                new_width = new_height * aspect_ratio
                img_shape.width = Emu(new_width)
                img_shape.height = Emu(new_height)
                
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_shape.left = slide_width - img_shape.width
                img_shape.top = slide_height - img_shape.height
                
                os.remove(temp_image_path.name)
                
        except Exception as e:
            st.write(f"Could not add image '{image}': {e}")

    prs.save(output_path)
    st.success("Presentation created successfully!")

st.title("PowerPoint Image Formatter")

# ZIP file uploader
uploaded_zip = st.file_uploader("Upload a ZIP file containing images", type=["zip"])
output_file_name = st.text_input("Specify Output PowerPoint File Name", value="output_presentation.pptx")
crop_left = st.number_input("Crop Pixels from Left", min_value=0, max_value=1000, value=250, step=10)
crop_bottom = st.number_input("Crop Pixels from Bottom", min_value=0, max_value=1000, value=42, step=10)

# Generate presentation button
if st.button("Generate Presentation"):
    if not uploaded_zip or not output_file_name:
        st.error("Please upload a ZIP file and specify the output file name.")
    else:
        with tempfile.TemporaryDirectory() as tmpdirname:
            # Extract and gather all image paths from the ZIP file
            images = extract_images_from_zip(uploaded_zip, tmpdirname)

            # Verify if images were found
            if not images:
                st.error("No valid images found in the uploaded ZIP file.")
            else:
                output_path = os.path.join(tmpdirname, output_file_name)
                
                # Run the presentation creation function
                create_widescreen_presentation_with_images(output_path, images, crop_left, crop_bottom)

                # Provide download link for the created presentation
                with open(output_path, "rb") as file:
                    st.download_button(
                        label="Download Presentation",
                        data=file,
                        file_name=output_file_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
'''

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image
from zipfile import ZipFile
import os
import tempfile

def extract_images_from_zip(zip_file, extract_to):
    image_extensions = {".jpg", ".jpeg", ".png", ".bmp", ".gif"}
    image_paths = []
    with ZipFile(zip_file) as zip_ref:
        # Extract all files and gather only valid image files, excluding macOS hidden files
        for file_info in zip_ref.infolist():
            if not file_info.is_dir():
                # Check for valid image extensions and exclude macOS hidden files
                if (os.path.splitext(file_info.filename)[1].lower() in image_extensions and 
                    not os.path.basename(file_info.filename).startswith("._")):
                    extracted_path = zip_ref.extract(file_info, path=extract_to)
                    image_paths.append(extracted_path)
    return image_paths

def create_widescreen_presentation_with_images(output_path, images, crop_left=0, crop_bottom=0, new_height_inches=6.0):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    if not images:
        st.warning("No images found to add to the presentation.")
        return

    for image in images:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        try:
            with Image.open(image) as img:
                original_width, original_height = img.size
                crop_box = (crop_left, 0, original_width, original_height - crop_bottom)
                cropped_img = img.crop(crop_box)
                
                temp_image_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                cropped_img.save(temp_image_path.name)
                
                img_shape = slide.shapes.add_picture(temp_image_path.name, 0, 0)
                aspect_ratio = img_shape.width / img_shape.height
                new_height = Inches(new_height_inches)
                new_width = new_height * aspect_ratio
                img_shape.width = Emu(new_width)
                img_shape.height = Emu(new_height)
                
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_shape.left = slide_width - img_shape.width
                img_shape.top = slide_height - img_shape.height
                
                os.remove(temp_image_path.name)
                
        except Exception as e:
            st.write(f"Could not add image '{image}': {e}")

    prs.save(output_path)
    st.success("Presentation created successfully!")

st.title("PowerPoint Image Formatter")

# ZIP file uploader
uploaded_zip = st.file_uploader("Upload a ZIP file containing images", type=["zip"])
output_file_name = st.text_input("Specify Output PowerPoint File Name", value="output_presentation.pptx")
crop_left = st.number_input("Crop Pixels from Left", min_value=0, max_value=1000, value=250, step=10)
crop_bottom = st.number_input("Crop Pixels from Bottom", min_value=0, max_value=1000, value=42, step=10)
new_height_inches = st.number_input("New Height in Inches", min_value=1.0, max_value=10.0, value=6.0, step=0.5)

# Generate presentation button
if st.button("Generate Presentation"):
    if not uploaded_zip or not output_file_name:
        st.error("Please upload a ZIP file and specify the output file name.")
    else:
        with tempfile.TemporaryDirectory() as tmpdirname:
            # Extract and gather all image paths from the ZIP file
            images = extract_images_from_zip(uploaded_zip, tmpdirname)

            # Verify if images were found
            if not images:
                st.error("No valid images found in the uploaded ZIP file.")
            else:
                output_path = os.path.join(tmpdirname, output_file_name)
                
                # Run the presentation creation function with user-defined height
                create_widescreen_presentation_with_images(output_path, images, crop_left, crop_bottom, new_height_inches)

                # Provide download link for the created presentation
                with open(output_path, "rb") as file:
                    st.download_button(
                        label="Download Presentation",
                        data=file,
                        file_name=output_file_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
